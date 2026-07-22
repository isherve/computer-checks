"""
Computer Checks – dark-mode thesis deck (12 slides)
No UI screenshots. Clean dark design with subtle computer/check atmosphere.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFilter, ImageEnhance

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "Computer_Checks_Thesis_Presentation_DARK_CLEAN.pptx"
ASSETS = ROOT / "Presentation_Assets"
ASSETS.mkdir(exist_ok=True)

NAVY = RGBColor(0x0B, 0x12, 0x1A)
TEAL = RGBColor(0x1A, 0xC4, 0xC8)
TEAL_DIM = RGBColor(0x0D, 0x73, 0x77)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
CARD = RGBColor(0x14, 0x1E, 0x2B)
CARD2 = RGBColor(0x1A, 0x27, 0x38)
LINE = RGBColor(0x2A, 0x3A, 0x4E)
GOLD = RGBColor(0xE8, 0xB8, 0x4A)

W = Inches(13.333)
H = Inches(7.5)
TOTAL = 12


def make_dark_bgs():
    """Abstract dark computer / check / QR atmosphere (no screenshots)."""
    w, h = 1920, 1080

    def canvas():
        img = Image.new("RGB", (w, h), (11, 18, 26))
        draw = ImageDraw.Draw(img, "RGBA")
        return img, draw

    # Main dark bg
    img, draw = canvas()
    draw.rectangle([0, 0, 40, h], fill=(13, 115, 119))
    for x in range(80, w, 80):
        draw.line([(x, 0), (x, h)], fill=(26, 196, 200, 16), width=1)
    for y in range(60, h, 80):
        draw.line([(0, y), (w, y)], fill=(26, 196, 200, 16), width=1)
    # QR-like motif top-right
    for i in range(8):
        for j in range(8):
            if (i + j) % 2 == 0 or (i * j) % 3 == 0:
                x0, y0 = 1580 + i * 32, 80 + j * 32
                draw.rectangle([x0, y0, x0 + 24, y0 + 24], fill=(13, 115, 119, 55))
    # check mark
    draw.line([(1500, 880), (1585, 955), (1760, 780)], fill=(26, 196, 200, 70), width=18)
    # laptop outline
    draw.rounded_rectangle([80, 820, 320, 980], radius=14, outline=(42, 58, 78, 120), width=3)
    draw.rectangle([110, 980, 290, 1000], fill=(42, 58, 78, 90))
    path_main = ASSETS / "bg_dark_main.png"
    img.save(path_main)

    # Title/closing denser
    img, draw = canvas()
    draw.rectangle([0, 0, 48, h], fill=(13, 115, 119))
    for x in range(100, w, 70):
        draw.line([(x, 0), (x, h)], fill=(20, 140, 145, 22), width=1)
    for i in range(7):
        for j in range(7):
            if (i + j) % 2 == 0:
                x0, y0 = 1600 + i * 30, 120 + j * 30
                draw.rectangle([x0, y0, x0 + 22, y0 + 22], fill=(26, 196, 200, 50))
    draw.line([(1480, 900), (1570, 980), (1760, 800)], fill=(26, 196, 200, 85), width=20)
    path_hero = ASSETS / "bg_dark_hero.png"
    img.save(path_hero)

    # Soft photo atmosphere from laptop image if available (not a UI screenshot)
    photo = ROOT / "QR" / "img" / "computer-tracking.jpg"
    path_photo = ASSETS / "bg_dark_photo.jpg"
    if photo.is_file():
        im = Image.open(photo).convert("RGB").resize((1920, 1080), Image.Resampling.LANCZOS)
        im = im.filter(ImageFilter.GaussianBlur(2.2))
        im = ImageEnhance.Brightness(im).enhance(0.22)
        im = ImageEnhance.Contrast(im).enhance(1.15)
        im = ImageEnhance.Color(im).enhance(0.7)
        overlay = Image.new("RGBA", im.size, (8, 14, 22, 180))
        mixed = Image.alpha_composite(im.convert("RGBA"), overlay).convert("RGB")
        d = ImageDraw.Draw(mixed)
        d.rectangle([0, 0, 40, 1080], fill=(13, 115, 119))
        mixed.save(path_photo, quality=92)
    else:
        path_photo = path_main

    return path_main, path_hero, path_photo


def set_run(run, size=18, bold=False, color=WHITE, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def put_bg(slide, path):
    pic = slide.shapes.add_picture(str(path), 0, 0, width=W, height=H)
    spTree = slide.shapes._spTree
    sp = pic._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def panel(slide, x, y, w, h, fill=CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    return shape


def footer(slide, page):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(0.42), W, Inches(0.42))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x07, 0x0C, 0x12)
    bar.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.45), H - Inches(0.36), Inches(10), Inches(0.28))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "Computer Checks  ·  QR Laptop Gate Verification  ·  UTB Rubavu Campus"
    set_run(r, 11, False, MUTED)
    num = slide.shapes.add_textbox(W - Inches(1.5), H - Inches(0.36), Inches(1.1), Inches(0.28))
    p = num.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{page} / {TOTAL}"
    set_run(r, 11, False, TEAL)


def title_block(slide, title, subtitle=None):
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.15))
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(0x08, 0x0E, 0x16)
    strip.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.15), W, Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL_DIM
    accent.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12.2), Inches(0.5))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = title
    set_run(r, 28, True, WHITE)
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12.2), Inches(0.3))
        r = s.text_frame.paragraphs[0].add_run()
        r.text = subtitle
        set_run(r, 13, False, MUTED)


def bullets(slide, items, left, top, width, size=14):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        r = p.add_run()
        r.text = "▸  " + item
        set_run(r, size, False, WHITE)


def build():
    bg_main, bg_hero, bg_photo = make_dark_bgs()
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    page = 0

    def n():
        nonlocal page
        page += 1
        return page

    # 1 Title
    s = prs.slides.add_slide(blank)
    put_bg(s, bg_hero)
    veil = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.4), Inches(9.5), Inches(4.5))
    veil.fill.solid()
    veil.fill.fore_color.rgb = RGBColor(0x08, 0x0E, 0x16)
    veil.line.color.rgb = LINE

    box = s.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(8.9), Inches(0.35))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "UTB RUBAVU CAMPUS  ·  FINAL YEAR PROJECT DEFENSE"
    set_run(r, 12, True, TEAL)

    box = s.shapes.add_textbox(Inches(0.9), Inches(2.25), Inches(8.9), Inches(1.3))
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "Computer Checks"
    set_run(r, 44, True, WHITE)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "QR-Based Laptop Gate Verification System"
    set_run(r, 20, False, MUTED)

    box = s.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(8.9), Inches(1.5))
    tf = box.text_frame
    for i, line in enumerate([
        "Presented by: Mfitumukiza Eric  &  Tuyishimire Gaspard",
        "Development model: Agile (iterative)",
        "Academic Year 2024 / 2025",
    ]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = para.add_run()
        r.text = line
        set_run(r, 15, False, WHITE)
        para.space_after = Pt(5)
    footer(s, n())

    # 2 Outline + Intro
    s = prs.slides.add_slide(blank)
    put_bg(s, bg_photo)
    title_block(s, "Outline & Introduction", "Defense roadmap and project context")
    panel(s, 0.45, 1.45, 6.1, 5.2)
    t = s.shapes.add_textbox(Inches(0.7), Inches(1.65), Inches(5.6), Inches(0.35))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "Presentation outline"
    set_run(r, 16, True, TEAL)
    bullets(s, [
        "Problem, objectives, Agile model",
        "Functional & non-functional requirements",
        "Analysis, design, architecture",
        "Implementation & gate workflow",
        "Testing, results, recommendations",
    ], 0.7, 2.2, 5.6, 15)

    panel(s, 6.8, 1.45, 6.0, 5.2)
    t = s.shapes.add_textbox(Inches(7.05), Inches(1.65), Inches(5.5), Inches(0.35))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "Introduction"
    set_run(r, 16, True, TEAL)
    bullets(s, [
        "Campus gates must verify personal computers",
        "Paper registers are slow and hard to audit",
        "QR codes enable fast phone-based logging",
        "One web platform for Admin and gate officers",
        "Built for UTB Rubavu Campus operations",
    ], 7.05, 2.2, 5.5, 14)
    footer(s, n())

    # 3 Problem + Objectives
    s = prs.slides.add_slide(blank)
    put_bg(s, bg_main)
    title_block(s, "Problem Statement & Objectives")
    panel(s, 0.45, 1.4, 6.1, 5.25)
    t = s.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.35))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "Problem"
    set_run(r, 16, True, GOLD)
    bullets(s, [
        "Manual gate logging causes delays and errors",
        "Illegible or incomplete paper records",
        "Difficult retrieval of past movements",
        "Weak digital ownership traceability",
    ], 0.7, 2.1, 5.6, 15)

    panel(s, 6.8, 1.4, 6.0, 5.25)
    t = s.shapes.add_textbox(Inches(7.05), Inches(1.55), Inches(5.5), Inches(0.35))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "Objectives"
    set_run(r, 16, True, TEAL)
    bullets(s, [
        "General: Design and implement a QR-based gate verification system for campus PCs",
        "Register laptops with owner details",
        "Generate printable QR codes",
        "Log check-in / check-out via scan",
        "Role-based access and PDF/CSV reports",
    ], 7.05, 2.1, 5.5, 14)
    footer(s, n())

    # 4 Agile
    s = prs.slides.add_slide(blank)
    put_bg(s, bg_main)
    title_block(s, "Model Used: Agile", "Iterative delivery of Computer Checks features")
    phases = [
        ("1", "Backlog", "Requirements &\nobjectives"),
        ("2", "Sprint plan", "Prioritize\nfeatures"),
        ("3", "Build", "PHP modules\n& database"),
        ("4", "Test", "Login, QR,\nlogs, reports"),
        ("5", "Review", "Demo &\nfeedback"),
    ]
    for i, (num, title, desc) in enumerate(phases):
        x = 0.4 + i * 2.55
        panel(s, x, 1.5, 2.4, 3.35, CARD2)
        oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.9), Inches(1.75), Inches(0.55), Inches(0.55))
        oval.fill.solid()
        oval.fill.fore_color.rgb = TEAL_DIM
        oval.line.fill.background()
        nt = s.shapes.add_textbox(Inches(x + 0.9), Inches(1.82), Inches(0.55), Inches(0.4))
        para = nt.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        r = para.add_run()
        r.text = num
        set_run(r, 15, True, WHITE)
        tt = s.shapes.add_textbox(Inches(x + 0.15), Inches(2.5), Inches(2.1), Inches(0.45))
        para = tt.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        r = para.add_run()
        r.text = title
        set_run(r, 15, True, TEAL)
        dd = s.shapes.add_textbox(Inches(x + 0.15), Inches(3.15), Inches(2.1), Inches(1.2))
        dd.text_frame.word_wrap = True
        para = dd.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        r = para.add_run()
        r.text = desc
        set_run(r, 12, False, MUTED)

    panel(s, 0.45, 5.1, 12.4, 1.4)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(5.3), Inches(12), Inches(1.0))
    tb.text_frame.word_wrap = True
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "Why Agile: registration → QR → scan → logs → reports were built and validated in short cycles. Thesis chapters follow a clear academic sequence (analysis → design → implementation → testing)."
    set_run(r, 14, False, WHITE)
    footer(s, n())

    # 5 FR
    s = prs.slides.add_slide(blank)
    put_bg(s, bg_main)
    title_block(s, "Functional Requirements", "What Computer Checks must do")
    frs = [
        ("FR1", "Authenticate Admin and Guest roles"),
        ("FR2", "Register laptop (SN, model, owner, ID)"),
        ("FR3", "Search PCs and generate QR codes"),
        ("FR4", "Scan QR → gate form → check-in/out"),
        ("FR5", "Save optional comments on each log"),
        ("FR6", "View and filter gate logs"),
        ("FR7", "Download reports as PDF / CSV"),
        ("FR8", "Admin user management"),
        ("FR9", "Change password securely"),
        ("FR10", "Show success message after commit"),
    ]
    for i, (code, text) in enumerate(frs):
        col = i % 2
        row = i // 2
        x = 0.45 + col * 6.4
        y = 1.4 + row * 0.9
        panel(s, x, y, 6.2, 0.8)
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.1), Inches(0.8))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = TEAL_DIM
        stripe.line.fill.background()
        tb = s.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.22), Inches(5.8), Inches(0.4))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = f"{code}   {text}"
        set_run(r, 14, False, WHITE)
    footer(s, n())

    # 6 NFR
    s = prs.slides.add_slide(blank)
    put_bg(s, bg_main)
    title_block(s, "Non-Functional Requirements", "Quality attributes of the system")
    nfrs = [
        ("Usability", "Clear dashboards; mobile gate form for phone scanning"),
        ("Performance", "Interactive response for login, QR, and log submit"),
        ("Security", "Password hashing, sessions, role-based page access"),
        ("Reliability", "Validated inputs; consistent auditable log records"),
        ("Maintainability", "Modular PHP: auth, registry, QR, logs, exports"),
        ("Portability", "XAMPP locally; PHP/MySQL hosting for production"),
        ("Availability", "Requires network for live QR scan URLs"),
        ("Scalability", "Database grows with computers and daily gate events"),
    ]
    for i, (title, desc) in enumerate(nfrs):
        col = i % 2
        row = i // 2
        x = 0.45 + col * 6.4
        y = 1.4 + row * 1.25
        panel(s, x, y, 6.2, 1.12)
        tt = s.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.18), Inches(5.7), Inches(0.3))
        r = tt.text_frame.paragraphs[0].add_run()
        r.text = title
        set_run(r, 15, True, TEAL)
        dd = s.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.55), Inches(5.7), Inches(0.45))
        dd.text_frame.word_wrap = True
        r = dd.text_frame.paragraphs[0].add_run()
        r.text = desc
        set_run(r, 13, False, WHITE)
    footer(s, n())

    # 7 Analysis & Design
    s = prs.slides.add_slide(blank)
    put_bg(s, bg_photo)
    title_block(s, "System Analysis & Design", "Actors, processes, and data stores")
    cols = [
        ("Actors", ["Admin", "Guest / gate officer", "Phone user (QR scan)"]),
        ("Main processes", ["Authenticate", "Register computer", "Generate QR", "Record gate log", "View / export logs"]),
        ("Core tables", ["users", "computer_info", "logs", "Linked by owner ID / SN"]),
    ]
    for i, (title, items) in enumerate(cols):
        x = 0.45 + i * 4.25
        panel(s, x, 1.45, 4.05, 5.2)
        t = s.shapes.add_textbox(Inches(x + 0.25), Inches(1.65), Inches(3.6), Inches(0.4))
        r = t.text_frame.paragraphs[0].add_run()
        r.text = title
        set_run(r, 16, True, TEAL)
        bullets(s, items, x + 0.25, 2.25, 3.6, 14)
    footer(s, n())

    # 8 Architecture
    s = prs.slides.add_slide(blank)
    put_bg(s, bg_main)
    title_block(s, "Architecture & Technology Stack", "Three-tier Computer Checks platform")
    tiers = [
        ("Presentation layer", "HTML, CSS, Bootstrap, JavaScript — dashboards, QR page, mobile gate form"),
        ("Application layer", "PHP + PDO + sessions — auth, validation, QR encoding, logging, PDF/CSV export"),
        ("Data layer", "MySQL (XAMPP / production hosting) — users, computer_info, logs"),
    ]
    for i, (title, body) in enumerate(tiers):
        y = 1.5 + i * 1.55
        panel(s, 0.5, y, 12.3, 1.4)
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.12), Inches(1.4))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = TEAL_DIM
        stripe.line.fill.background()
        tt = s.shapes.add_textbox(Inches(0.9), Inches(y + 0.25), Inches(11.6), Inches(0.35))
        r = tt.text_frame.paragraphs[0].add_run()
        r.text = title
        set_run(r, 18, True, TEAL)
        bb = s.shapes.add_textbox(Inches(0.9), Inches(y + 0.7), Inches(11.6), Inches(0.5))
        bb.text_frame.word_wrap = True
        r = bb.text_frame.paragraphs[0].add_run()
        r.text = body
        set_run(r, 14, False, WHITE)
    footer(s, n())

    # 9 Implementation
    s = prs.slides.add_slide(blank)
    put_bg(s, bg_main)
    title_block(s, "Implementation Highlights", "What was delivered in Computer Checks")
    features = [
        ("Register & QR", "Create laptop records and print scannable QR codes"),
        ("Gate scan flow", "Phone opens form → Check-In/Out → success message"),
        ("Logs hub", "Activity table, comments, period filters"),
        ("Exports", "Download filtered logs as PDF or CSV"),
        ("Roles", "Admin manages users; Guest runs gate operations"),
        ("UTB IDs", "Student registration numbers use UTB format"),
    ]
    for i, (title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        x = 0.45 + col * 4.2
        y = 1.45 + row * 2.4
        panel(s, x, y, 4.0, 2.2)
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.12), Inches(2.2))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = TEAL_DIM
        stripe.line.fill.background()
        tt = s.shapes.add_textbox(Inches(x + 0.35), Inches(y + 0.45), Inches(3.4), Inches(0.45))
        r = tt.text_frame.paragraphs[0].add_run()
        r.text = title
        set_run(r, 17, True, TEAL)
        dd = s.shapes.add_textbox(Inches(x + 0.35), Inches(y + 1.05), Inches(3.4), Inches(0.8))
        dd.text_frame.word_wrap = True
        r = dd.text_frame.paragraphs[0].add_run()
        r.text = desc
        set_run(r, 13, False, MUTED)
    footer(s, n())

    # 10 Workflow
    s = prs.slides.add_slide(blank)
    put_bg(s, bg_photo)
    title_block(s, "Operational Workflow", "From registration to auditable gate log")
    steps = ["Register\nlaptop", "Generate\nQR", "Print /\nkeep QR", "Scan at\ngate", "Submit\nIn / Out", "Log &\nreport"]
    for i, label in enumerate(steps):
        x = 0.4 + i * 2.15
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(1.95), Inches(1.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = TEAL_DIM if i % 2 == 0 else CARD2
        shape.line.color.rgb = LINE
        tb = s.shapes.add_textbox(Inches(x + 0.1), Inches(2.5), Inches(1.75), Inches(1.2))
        tb.text_frame.word_wrap = True
        para = tb.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        r = para.add_run()
        r.text = f"{i + 1}\n{label}"
        set_run(r, 13, True, WHITE)

    panel(s, 0.45, 4.4, 12.4, 2.15)
    tb = s.shapes.add_textbox(Inches(0.75), Inches(4.65), Inches(11.9), Inches(1.6))
    tb.text_frame.word_wrap = True
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "After commit, the phone shows: “A log for [serial] whose the owner is [name] is recorded successfully! Status: [check-in/check-out].” Comments entered at the gate appear under Logs and in PDF/CSV exports."
    set_run(r, 15, False, WHITE)
    footer(s, n())

    # 11 Testing + Conclusion
    s = prs.slides.add_slide(blank)
    put_bg(s, bg_main)
    title_block(s, "Testing, Results & Conclusion")
    panel(s, 0.45, 1.4, 6.1, 5.25)
    t = s.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.35))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "Testing & results"
    set_run(r, 16, True, GOLD)
    bullets(s, [
        "Admin / Guest login validated",
        "Laptop registration & UTB IDs OK",
        "QR opens gate form on phone",
        "Check-in/out + comments stored",
        "PDF & CSV exports verified",
        "Usable on desktop and mobile",
    ], 0.7, 2.1, 5.6, 14)

    panel(s, 6.8, 1.4, 6.05, 5.25)
    t = s.shapes.add_textbox(Inches(7.05), Inches(1.55), Inches(5.6), Inches(0.35))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "Conclusion & recommendations"
    set_run(r, 16, True, TEAL)
    bullets(s, [
        "Computer Checks digitizes PC gate control",
        "Faster logging and stronger audit trails",
        "Train officers; keep QR on each device",
        "Use durable MySQL for production hosting",
        "Future: alerts, offline mode, more assets",
    ], 7.05, 2.1, 5.6, 14)
    footer(s, n())

    # 12 Thank you
    s = prs.slides.add_slide(blank)
    put_bg(s, bg_hero)
    veil = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.6), Inches(12.2), Inches(4.5))
    veil.fill.solid()
    veil.fill.fore_color.rgb = RGBColor(0x08, 0x0E, 0x16)
    veil.line.color.rgb = LINE

    t = s.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(11.5), Inches(0.8))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "Thank You"
    set_run(r, 48, True, WHITE)

    t = s.shapes.add_textbox(Inches(0.95), Inches(2.9), Inches(11.5), Inches(0.4))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "Questions from the Panel"
    set_run(r, 22, False, TEAL)

    t = s.shapes.add_textbox(Inches(0.95), Inches(3.6), Inches(11.5), Inches(2.0))
    tf = t.text_frame
    for i, line in enumerate([
        "Mfitumukiza Eric  &  Tuyishimire Gaspard",
        "Computer Checks · UTB Rubavu Campus",
        "Agile model  ·  PHP / MySQL / QR",
        "Live demo: https://computer-checks.vercel.app",
    ]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = para.add_run()
        r.text = line
        set_run(r, 15, False, MUTED if i else WHITE)
        para.space_after = Pt(5)
    footer(s, n())

    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
